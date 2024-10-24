import os
import shutil

from pptx import Presentation

def extract_text_from_pptx(pptx_file):
    # Load the presentation
    presentation = Presentation(pptx_file)
    
    # Initialize an empty list to store slide texts
    slides_text = []
    
    # Loop through each slide
    for slide in presentation.slides:
        slide_text = []
        # Loop through all shapes in the slide
        for shape in slide.shapes:
            # If the shape has text, extract it
            if hasattr(shape, "text"):
                slide_text.append(shape.text)
        
        # Join all text in the slide and add it to slides_text list
        slides_text.append("\n".join(slide_text))
    
    # Join text from all slides
    full_text = "\n\n".join(slides_text)
    
    return full_text

if __name__ == '__main__':
    save_path = './data/documents'
    pptx_path = './data/pptx'

    if os.path.exists(save_path):
        shutil.rmtree(save_path)
    os.mkdir(save_path)

    for pptx in os.listdir(pptx_path):
        pptx_file = os.path.join(pptx_path,pptx)
        text = extract_text_from_pptx(pptx_file)
        with open(os.path.join(save_path,pptx), 'w') as f:
            f.write(text)





